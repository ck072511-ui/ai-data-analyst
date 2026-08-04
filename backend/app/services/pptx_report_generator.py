import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from typing import Any, Dict
from pptx import Presentation
from pptx.util import Inches, Pt

class PPTXReportGenerator:
    @staticmethod
    def generate(file_path: str, data: Dict[str, Any], branding: Dict[str, Any]):
        """Generates a professional PowerPoint presentation using python-pptx."""
        prs = Presentation()

        company_name = branding.get("company_name", "Enterprise Systems")
        version = branding.get("version", "1.0.0")

        # Slide 1: Title Slide
        slide_layout = prs.slide_layouts[0] # Title Layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = data.get("title", "Enterprise Analytics Report")
        slide.placeholders[1].text = f"Compiled for: {company_name}\nVersion: {version} | Generated: {data.get('timestamp')}"

        # Slide 2: Executive Summary
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Executive Summary"
        slide.placeholders[1].text = data.get("final_answer", "Analyzed dataset metrics successfully.")

        # Slide 3: KPIs
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Key Performance Indicators"
        tf = slide.placeholders[1].text_frame
        kpis = data.get("kpis", [])
        if kpis:
            for k in kpis[:4]:
                p = tf.add_paragraph()
                p.text = f"• {k.get('title')}: {k.get('value')}"
                p.level = 0
        else:
            tf.text = "No KPI elements recommended."

        # Slide 4: Chart slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Performance Metric Visualization"
        chart_path = f"{file_path}.png"
        try:
            plt.figure(figsize=(6, 3.5))
            plt.plot([1, 2, 3, 4], [10, 20, 25, 30], label="Sales Value", marker='d', color='#ea580c')
            plt.title("Sample Performance Metrics Trend")
            plt.grid(True)
            plt.tight_layout()
            plt.savefig(chart_path)
            plt.close()
            
            # Insert image flowable on slide
            slide.shapes.add_picture(chart_path, Inches(1.5), Inches(2.0), width=Inches(7.0))
        except Exception as e:
            slide.placeholders[1].text = f"Failed to generate trend graphic: {e}"

        # Slide 5: AI Insights & Risks
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "AI Business Insights & Risks"
        tf = slide.placeholders[1].text_frame
        
        insights = data.get("insights", [])
        tf.text = "Recommended Actions:"
        for ins in insights[:2]:
            p = tf.add_paragraph()
            p.text = f"• {ins}"
            p.level = 1

        p_risk = tf.add_paragraph()
        p_risk.text = "Identified Risk Vectors:"
        p_risk.level = 0
        for r in data.get("trends", ["Dataset complexity thresholds"]):
            p = tf.add_paragraph()
            p.text = f"• {r}"
            p.level = 1

        # Save presentation
        prs.save(file_path)

        # Cleanup image
        if os.path.exists(chart_path):
            try:
                os.remove(chart_path)
            except:
                pass
